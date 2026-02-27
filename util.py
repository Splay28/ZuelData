from Crypto.Cipher import AES
import base64
from PIL import Image, ImageDraw, ImageFont
from random import choice, randint, randrange
import string
from io import BytesIO
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication
from email.header import Header

import re
import PyPDF2
import xlrd

# REF https://stackoverflow.com/questions/55497789/find-a-word-in-multiple-powerpoint-files-python/55763992#55763992
import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from docx import Document

# 候选字符集,大小写字母+数字
chrs = string.ascii_letters + string.digits
aes_key = 'k_aes'
BLOCK_SIZE = 16  # Bytes
pad = lambda s: s + (BLOCK_SIZE - len(s) % BLOCK_SIZE) * \
                chr(BLOCK_SIZE - len(s) % BLOCK_SIZE)
unpad = lambda s: s[:-ord(s[len(s) - 1:])]


def aesEncrypt(key, data, salt):
    '''
    AES的ECB模式加密方法
    :param key: 密钥
    :param data:被加密字符串（明文）
    :return:密文
    '''
    data += salt
    key = key.encode('utf8')
    # 字符串补位
    data = pad(data)
    cipher = AES.new(key, AES.MODE_ECB)
    # 加密后得到的是bytes类型的数据，使用Base64进行编码,返回byte字符串
    result = cipher.encrypt(data.encode())
    encodestrs = base64.b64encode(result)
    enctext = encodestrs.decode('utf8')
    return enctext

def aesDecrypt(key, data, salt):
    '''

    :param key: 密钥
    :param data: 加密后的数据（密文）
    :return:明文
    '''
    key = key.encode('utf8')
    data = base64.b64decode(data)
    cipher = AES.new(key, AES.MODE_ECB)

    # 去补位
    text_decrypted = unpad(cipher.decrypt(data))
    text_decrypted = text_decrypted.decode('utf8')
    text_decrypted = text_decrypted[0:-(len(salt))]
    return text_decrypted



def selected_chrs(length):
  """
  返回length个随机字符串
  :param length:
  :return:
  """
  result = ''.join(choice(chrs) for _ in range(length))
  return result

def get_color():
  """
  设置随机颜色
  :return:
  """
  r = randint(0, 255)
  g = randint(0, 255)
  b = randint(0, 255)
  return (r, g, b)

def getcode(size=(200, 100), chrNumber=4, bgcolor=(255, 255, 255)):
  bytesIO = BytesIO()
  """
  定义图片大小，验证码长度，背景颜色
  :param size:
  :param chrNumber:
  :param bgcolor:
  :return:元组 0为图片，1为码
  """
  # 创建空白图像和绘图对象
  image_tmp = Image.new('RGB', size, bgcolor)
  draw = ImageDraw.Draw(image_tmp)

  # 生成并计算随机字符的宽度和高度
  text = selected_chrs(chrNumber)
  font = ImageFont.truetype('./static/arial.ttf', 48)
  width, height = draw.textsize(text, font)
  if width + 2*chrNumber > size[0] or height > size[1]:
    return

  # 绘制字符串
  startX = 0
  width_eachchr = width // chrNumber # 计算每个字符宽度
  for i in range(chrNumber):
    startX += width_eachchr + 1
    position = (startX, (size[1]-height)//2+randint(-10, 10)) # 字符坐标, Y坐标上下浮动
    draw.text(xy=position, text=text[i], font=font, fill=get_color()) # 绘制函数

  # 对像素位置进行微调，实现验证码扭曲效果
  img_final = Image.new('RGB', size, bgcolor)
  pixels_final = img_final.load()
  pixels_tmp = image_tmp.load()
  for y in range(size[1]):
    offset = randint(-1, 0) # randint()相当于闭区间[x,y]
    for x in range(size[0]):
      newx = x + offset # 像素微调
      if newx >= size[0]:
        newx = size[0] - 1
      elif newx < 0:
        newx = 0
      pixels_final[newx, y] = pixels_tmp[x, y]

  # 绘制随机颜色随机位置的干扰像素
  draw = ImageDraw.Draw(img_final)
  for i in range(int(size[0]*size[1]*0.07)): # 7%密度的干扰像素
    draw.point((randrange(size[0]), randrange(size[1])), fill=get_color()) # randrange取值范围是左开右闭

  # 绘制随机干扰线，这里设置为8条
  for i in range(8):
    start = (0, randrange(size[1]))
    end = (size[0], randrange(size[1]))
    draw.line([start, end], fill=get_color(), width=1)

  # 绘制随机弧线
  for i in range(8):
    start = (-50, -50) # 起始位置在外边看起来才会像弧线
    end = (size[0]+10, randint(0, size[1]+10))
    draw.arc(start+end, 0, 360, fill=get_color())

  # 保存图片
  img_final.save(bytesIO, format='PNG')
  data = bytesIO.getvalue()
  data = base64.b64encode(data)
  data = data.decode('utf-8')

  
  return (data,text)

# 第三方 SMTP 服务
mail_host="smtp.163.com"  #设置服务器
mail_user="_"    #用户名
mail_pass="_"   #口令 
sender = '_'

def email(receiver,title,abstract,files,task_url,subdate,free_u=0):
  """
  定义图片大小，验证码长度，背景颜色
  :param receiver:收件人邮箱
  :param title:标题
  :param abstract:内容
  :param files:附加文件路径列表
  :param task_url:任务链接
  :param subdate:截止日期
  :param free_u:默认为0，是否自由使用文本发送邮件，否则按任务格式发送邮件，为1留空task_url和subdate
  :return:正常0，否则返回错误
  """
  #创建一个带附件的实例
  message = MIMEMultipart('mixed')
  message['From'] = Header(u'中南法援 <%s>'%sender, 'utf-8')
  message['To'] =  Header(u'诉讼部 <%s>'%receiver, 'utf-8')
  message['Subject'] = Header(title, 'utf-8').encode()
  #邮件正文内容
  if(free_u):
    message.attach(MIMEText(abstract + '\n请勿回复此邮件！', 'plain', 'utf-8'))
  elif(subdate):
    #根据是否存在截止日期判断是否办结
    message.attach(MIMEText(abstract + '\n\n请在：' + str(subdate) + '之前在以下网页提交\n' + task_url + '\n请勿回复此邮件！', 'plain', 'utf-8'))
  else:
    message.attach(MIMEText(abstract + '\n\n位于' + task_url + ' 的案件现已办结，特此告知！' + '\n请勿回复此邮件！', 'plain', 'utf-8'))
  
  if(files):
    for i in files:
      #获取最新的email路径，如果存在说明有报错，构造附件，发送email路径下的excel文件
      att1 = MIMEApplication(open(i,'rb').read())
      att1['Content-Type'] = 'application/octet-stream'
      #这里的filename可以任意写，写什么邮件就显示什么名字
      filename=i.split('/')[-1:][0]
      att1.add_header('Content-Disposition', 'attachment', filename=filename)
      message.attach(att1)

  try:
    smtpObj = smtplib.SMTP() 
    smtpObj.connect(mail_host, 25)    # 25 为 SMTP 端口号
    smtpObj.login(mail_user,mail_pass)  
    smtpObj.sendmail(sender, receiver, message.as_string())
    return 0 
  except smtplib.SMTPException as e:
    return e



def inputpdf(s):
    l = []
    nums = re.findall(r'\d+', s)  
    chars = re.findall(r'\D+', s)
    l.append(re.findall('[^\x00-\xff]+',chars[0])[0])
    l.append(chars[0].replace(l[0], '') + nums[1])
    '''
    l.append(re.findall(r'通识教育课程|素质教育课程|专业教育课程', s)[0])
    l.append(re.findall(r'(.*?)(通识教育课程|素质教育课程|专业教育课程)', s)[0][0].replace(nums[0], '').replace(l[0], '').replace(l[1], ''))
    l.append('通识选修')
    l.append(nums[2][:1:])
    l.append(nums[2][1::])
    '''
    l.append(re.findall(r'全面发展|科学素养|家国情怀|文化传承|国际视野|经法管融通', s)[0])
    return l

def read_pdf(path = './t.pdf'):
    num = []
    module = []

    pdffile = open(path, 'rb')  # 读取pdf文件
    pdfreader = PyPDF2.PdfFileReader(pdffile)# 读取pdf页数总数
    pdfnum = pdfreader.numPages
    for i in range(pdfnum):
        page0 = pdfreader.getPage(i)
        templ = page0.extractText().replace(' ','').split('\n')
        for j in templ:
            if re.findall(r'学年第.学期', j) or re.findall(r'第.页', j) or re.findall(r'课程编号', j) or re.findall(r'一览表', j) or re.findall(r'课程名称', j):
                continue
            num.append(inputpdf(j)[1])
            module.append(inputpdf(j)[2])


    pdffile.close()
    return [num, module]

def read_xls(path = 't.xls'):
    coursel = []
    courses = xlrd.open_workbook(path)
    table = courses.sheets()[0]
    nrows = table.nrows

    std = ['课程号', '课序号', '课程名', '学时', '学分', '任课教师', '上课时间', '周学时', '上课周次', '选课限制说明', '上课班级', '上课年级', '课程性质']

    if(table.row_values(1) != std):
        return -1

    for i in range(nrows):
        if i <= 1:
            continue
        t = table.row_values(i)
        coursel.append(t)

    return coursel

def mix(path_x='./t.xls', path_p='./t.pdf'):
    lx = read_xls(path_x)
    if(lx==-1):return -1
    lp = read_pdf(path_p)
    result = []
    for i in lx:
        if(i[0] in lp[0]):
            m = lp[1][lp[0].index(i[0])]
            i.append(m)
        else:
            i.append(-1)
    
        if(not i[9]):i[9]='-1'

        result.append(i)

    return result

def read_arrange(path = 't.xls'):
    coursel = []
    courses = xlrd.open_workbook(path)
    table = courses.sheets()[0]
    nrows = table.nrows

    std = ['姓名', '邮箱', '年级']

    if(table.row_values(1) != std):
        return -1

    for i in range(nrows):
        if i <= 1:
            continue
        t = table.row_values(i)
        if(t):
           if(t[0]):
            coursel.append(t)

    return coursel

def read_lawsuit(path = 't.xls'):
    arrangel = []
    xls = xlrd.open_workbook(path)
    table = xls.sheets()[0]
    nrows = table.nrows

    std = ['时间', '姓名', '邮箱', '姓名', '邮箱']

    if(table.row_values(1) != std):
        return -1

    for i in range(nrows):
        if i <= 1:
            continue
        t = table.row_values(i)
        if(t):
           if(t[1] and t[2] and t[3] and t[4]):
            arrangel.append([t[1], t[2], t[3], t[4]])

    return arrangel

def read_horo(path = 't.xls'):
    coursel = []
    courses = xlrd.open_workbook(path)
    table = courses.sheets()[0]
    nrows = table.nrows

    std = ['一件事', '一个好的评价', '一个坏的评价']

    if(table.row_values(1) != std):
        return -1

    for i in range(nrows):
        if i <= 1:
            continue
        t = table.row_values(i)
        if(t):
           if(t[0]):
            coursel.append(t)

    return coursel

def read_quota(path = 't.xls'):
    coursel = []
    courses = xlrd.open_workbook(path)
    table = courses.sheets()[0]
    nrows = table.nrows

    std = ['出处', '内容']

    if(table.row_values(1) != std):
        return -1

    for i in range(nrows):
        if i <= 1:
            continue
        t = table.row_values(i)
        if(t):
           if(t[0]):
            coursel.append(t)

    return coursel



def get_sample(target, fulltext):
  sample = ''
  if target in fulltext:
      pos = fulltext.index(target)
      l = 0
      r = len(fulltext) - 1
      if(pos >= 7):
          l = pos - 7
          if(pos <= r -13):
              r = pos + 13
          else:
              r = pos + len(target)
              leng = 20 - len(target)
              if(pos - leng >= 0):
                  l = pos - leng
              else:
                  l = 0
      else:
          l = 0
          if(r >= 21):
              r = 21
          else:
              r = pos + len(target)
              
      sample = fulltext[l:r]
  return sample

def CheckRecursivelyForText(txt, shpthissetofshapes, filename):
  ts = []
  if(len(txt) > 10):txt=txt[:10]
  for shape in shpthissetofshapes:
      if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
          ts.extend(CheckRecursivelyForText(txt , shape.shapes, filename))
      else:
          if hasattr(shape, "text"):
              shape.text = shape.text.lower()
              shape.text = shape.text.replace('\n',' ')
              if txt in shape.text:
                  sample = get_sample(txt, shape.text)

                  ts.append(sample)
                  #break
  return ts

def find_pptx(txt, path = "./files/test/"):
  files = [x for x in os.listdir(path) if x.endswith(".pptx")]
  result = {}
  #文件名：样本列表
  for eachfile in files:
      sl = []
      prs = Presentation(path + eachfile) 
      for slide in prs.slides:
          tsl = CheckRecursivelyForText(txt ,slide.shapes, eachfile)
          sl.extend(tsl)
      if(sl):
          result[path+eachfile] = sl
  
  return result

def get_doc_path (path) :
  file_list = os.listdir(path)
  # 正则匹配路径下所有.docx结尾的文件
  doc_list = [i for i in file_list if re.compile(r'\w+.docx').match(i)]
  #拼接ur进入doc_list，获得完整路径
  for i in range(len(doc_list)):
      doc_list[i] = path + '/' + doc_list[i]
  return doc_list
#以段落为单位切片，并查找关键字
def find_text (path,word):
  document = Document(path)
  all_paragraphs = document.paragraphs
  list1 = []
  for paragraph in all_paragraphs:
      str1 = paragraph.text
      if str1.find(word) != -1 :
          list1.append(str1)
  return list1

def find_word (key_word, path = './files/test/'):
  doc_path = get_doc_path(path)
  findlist={}
  for x in doc_path :
      result = find_text(x,key_word)
      if(result):
          findlist[x]=result
  return findlist

def aggregate_search(txt, path = './files/test/'):
  t0 = find_pptx(txt, path)
  t1 = find_word(txt, path)

  result = [t0, t1]

  return result

def aggregate_search_list(txt, doclist, pptlist):
  findlist={}
  for x in doclist :
      result = find_text(x,txt)
      if(result):
          findlist[x]=result
  #文件名：样本列表
  for eachfile in pptlist:
      sl = []
      prs = Presentation(eachfile) 
      for slide in prs.slides:
          tsl = CheckRecursivelyForText(txt ,slide.shapes, eachfile)
          sl.extend(tsl)
      if(sl):
          findlist[eachfile] = sl

  return findlist
